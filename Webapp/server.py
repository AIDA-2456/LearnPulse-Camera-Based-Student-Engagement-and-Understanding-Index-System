import os, json, shutil, sqlite3, hashlib, secrets, threading, time, uuid
import subprocess, tempfile
import re
from collections import defaultdict, deque
from datetime import datetime

import numpy as np
import cv2
try:
    from dotenv import load_dotenv
    load_dotenv() 
except ImportError:
    pass
from fastapi import FastAPI, UploadFile, File, Form, HTTPException, Cookie, Response
from fastapi.responses import FileResponse, JSONResponse
import uvicorn

app = FastAPI(title="LearnPulse")

UPLOAD_DIR, LOG_DIR, DB_PATH = "Webapp/uploads", "Webapp/logs", "Webapp/learnpulse.db"
os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(LOG_DIR, exist_ok=True)

GRID_ROWS, GRID_COLS = 3, 4     

SAMPLES_PER_SEC = 1.5    
INFER_WIDTH     = 512     
LIVE_SAMPLES_PS = 2.0     

def db():
    con = sqlite3.connect(DB_PATH, check_same_thread=False)
    con.row_factory = sqlite3.Row
    return con

with db() as con:
    con.execute("""CREATE TABLE IF NOT EXISTS users(
        id TEXT PRIMARY KEY, name TEXT, email TEXT UNIQUE,
        salt TEXT, pwd TEXT, created TEXT)""")
    con.execute("""CREATE TABLE IF NOT EXISTS tokens(
        token TEXT PRIMARY KEY, user_id TEXT, created TEXT)""")
    con.execute("""CREATE TABLE IF NOT EXISTS sessions(
        id TEXT PRIMARY KEY, user_id TEXT, subject TEXT, class_name TEXT,
        topics TEXT, mode TEXT, created TEXT, duration_s INTEGER,
        average REAL, lowest_pct REAL, lowest_time INTEGER,
        coverage REAL, gap_seconds INTEGER, phone_peak INTEGER, phone_share REAL,
        events_json TEXT, series_json TEXT, zones_json TEXT)""")
    try:
        con.execute("ALTER TABLE sessions ADD COLUMN slides_json TEXT")
    except sqlite3.OperationalError:
        pass         


PW_RULES = [
    ("at least 8 characters", lambda p: len(p) >= 8),
    ("a lower-case letter", lambda p: any(c.islower() for c in p)),
    ("an upper-case letter", lambda p: any(c.isupper() for c in p)),
    ("a number", lambda p: any(c.isdigit() for c in p)),
    ("a symbol such as ! ? @ # or -",
     lambda p: any(not c.isalnum() and not c.isspace() for c in p)),
]


def check_password(pw):
    """Return a list of unmet requirements."""
    return [name for name, test in PW_RULES if not test(pw)]


def hash_pw(pw, salt):
    return hashlib.pbkdf2_hmac("sha256", pw.encode(), salt.encode(), 120000).hex()


def user_from_token(token):
    if not token:
        return None
    with db() as con:
        r = con.execute(
            "SELECT u.* FROM tokens t JOIN users u ON u.id=t.user_id WHERE t.token=?",
            (token,)).fetchone()
    return dict(r) if r else None


def require(token):
    u = user_from_token(token)
    if not u:
        raise HTTPException(401, "Please sign in again.")
    return u

from ultralytics import YOLO
print("Loading vision models...")
pose_model = YOLO("yolov8n-pose.pt")
object_model = YOLO("yolov8n.pt")

USE_MODEL = False
try:
    import joblib
    _b = joblib.load("engagement_model.joblib")
    clf, FEATURE_ORDER = _b["model"], _b["features"]
    USE_MODEL = True
    print("Trained engagement model loaded.")
except Exception as e:
    print("No trained model found - rules only:", e)

NOSE, L_EYE, R_EYE, L_EAR, R_EAR = 0, 1, 2, 3, 4
L_SH, R_SH = 5, 6
LOWER = [11, 12, 13, 14, 15, 16]
MIN_BOX_FRAC, MIN_FACE_CONF = 0.18, 0.55
PHONE_CLASS, PHONE_CONF, PHONE_EVERY = 67, 0.25, 6


def signals(person, box):
    x1, y1, x2, y2 = (float(v) for v in box)
    bh = max(y2 - y1, 1.0)
    n_, le, re = person[NOSE], person[L_EYE], person[R_EYE]
    la, ra = person[L_EAR], person[R_EAR]
    ls, rs = person[L_SH], person[R_SH]
    s = {"nose_conf": float(n_[2]), "eye_conf": (float(le[2]) + float(re[2])) / 2,
         "nose_x": float(n_[0]) / bh if n_[2] > .5 else None,
         "nose_y": float(n_[1]) / bh if n_[2] > .5 else None,
         "face_ratio": None, "head_down": None,
         "nose_to_shoulder": None, "shoulder_width": None, "eye_nose_dist": None}
    if n_[2] > .5 and la[2] > .3 and ra[2] > .3:
        lx, rx = float(min(la[0], ra[0])), float(max(la[0], ra[0]))
        if rx > lx:
            s["face_ratio"] = (float(n_[0]) - lx) / (rx - lx)
    if n_[2] > .5 and ls[2] > .3 and rs[2] > .3:
        sy = (float(ls[1]) + float(rs[1])) / 2
        s["head_down"] = 1.0 if float(n_[1]) >= sy else 0.0
        s["nose_to_shoulder"] = (sy - float(n_[1])) / bh
        s["shoulder_width"] = abs(float(rs[0]) - float(ls[0])) / bh
    if n_[2] > .5 and le[2] > .3 and re[2] > .3:
        ex, ey = (float(le[0]) + float(re[0])) / 2, (float(le[1]) + float(re[1])) / 2
        s["eye_nose_dist"] = float(np.hypot(float(n_[0]) - ex, float(n_[1]) - ey)) / bh
    return s


def features(hist):
    def col(k): return [f[k] for f in hist if f[k] is not None]
    def m(x): return float(np.mean(x)) if x else -1.0
    def sd(x): return float(np.std(x)) if len(x) > 1 else -1.0
    r_, xs, ys = col("face_ratio"), col("nose_x"), col("nose_y")
    if len(xs) > 2:
        st = np.hypot(np.diff(xs), np.diff(ys))
        mv, mx = float(np.mean(st)), float(np.max(st))
        big = int(np.sum(st > np.mean(st) + 2 * np.std(st))) if len(st) > 2 else 0
        th = np.median(st); longest = cur = 0
        for v in st:
            cur = cur + 1 if v <= th else 0
            longest = max(longest, cur)
        stf = float(np.mean(st <= th)); half = len(st) // 2
        drift = float(np.mean(st[half:]) - np.mean(st[:half])) if half else 0.0
    else:
        mv = mx = -1.0; big = longest = -1; stf = -1.0; drift = 0.0
    return {"face_ratio_mean": m(r_), "face_ratio_std": sd(r_),
            "facing_front_frac": (float(np.mean([1 if .3 < v < .7 else 0 for v in r_]))
                                  if r_ else -1.0),
            "head_down_frac": m(col("head_down")), "nose_conf_mean": m(col("nose_conf")),
            "nose_x_std": sd(xs), "nose_y_std": sd(ys), "eye_conf_mean": m(col("eye_conf")),
            "eye_nose_dist_mean": m(col("eye_nose_dist")),
            "eye_nose_dist_std": sd(col("eye_nose_dist")),
            "nose_to_shoulder_mean": m(col("nose_to_shoulder")),
            "nose_to_shoulder_std": sd(col("nose_to_shoulder")),
            "shoulder_width_mean": m(col("shoulder_width")),
            "shoulder_width_std": sd(col("shoulder_width")),
            "move_mean": mv, "move_max": mx, "big_moves": big,
            "longest_still": longest, "still_frac": stf, "movement_drift": drift}


def judge_frame(frame, history, phone_streak, cached_phones, frame_h, frame_w, zones):
    """Analyse one frame -> (engaged, disengaged, phones)."""
    eng = dis = phones = 0
    for result in pose_model.track(frame, persist=True, conf=.4, verbose=False,
                                   imgsz=INFER_WIDTH):
        if result.keypoints is None or len(result.keypoints.data) == 0:
            continue
        kp, boxes = result.keypoints.data, result.boxes.xyxy
        ids = result.boxes.id
        n = len(kp)
        t_idx, best = -1, -1
        for p in range(n):
            x1, y1, x2, y2 = (float(v) for v in boxes[p])
            sc = sum(1 for k in LOWER if float(kp[p][k][2]) > .3) * 1e6 + (x2-x1)*(y2-y1)
            if sc > best:
                best, t_idx = sc, p
        for p in range(n):
            if p == t_idx:
                continue
            x1, y1, x2, y2 = (float(v) for v in boxes[p])
            tid = int(ids[p]) if ids is not None else -1
            sig = signals(kp[p], boxes[p])
            history[tid].append(sig)
            has = any(x1-20 <= cx <= x2+20 and y1-20 <= cy <= y2+20
                      for (cx, cy) in cached_phones)
            phone_streak[tid] = phone_streak[tid] + 1 if has else 0
            if phone_streak[tid] >= 2:
                is_eng = False; phones += 1
            else:
                rec = list(history[tid])[-3:]
                hd = [f["head_down"] for f in rec if f["head_down"] is not None]
                hdn = (np.mean(hd) > .6) if hd else False
                if (USE_MODEL and (y2-y1)/frame_h >= MIN_BOX_FRAC
                        and (sig["nose_conf"]+sig["eye_conf"])/2 >= MIN_FACE_CONF
                        and not hdn and len(history[tid]) >= 6):
                    vec = np.array([[features(history[tid])[k] for k in FEATURE_ORDER]])
                    is_eng = bool(clf.predict(vec)[0] == 1)
                elif sig["head_down"] == 1.0:
                    is_eng = True
                elif sig["face_ratio"] is not None:
                    is_eng = .3 < sig["face_ratio"] < .7
                else:
                    is_eng = True
            eng += 1 if is_eng else 0
            dis += 0 if is_eng else 1
            cx, cy = (x1+x2)/2, (y1+y2)/2
            r = min(GRID_ROWS-1, int(cy / frame_h * GRID_ROWS))
            c = min(GRID_COLS-1, int(cx / frame_w * GRID_COLS))
            z = zones[r][c]
            z["total"] += 1
            z["eng"] += 1 if is_eng else 0
    return eng, dis, phones


def summarise(times, engs, phones, zones, duration):
    t_all = np.array(times, float); e_all = np.array(engs, float)
    p_all = np.array(phones, float)
    valid = e_all > 0
    coverage = float(valid.mean()*100) if len(e_all) else 0
    gaps = int((~valid).sum())
    if valid.sum() < 5:
        return {"error": "Students were visible for too little of this recording."}
    t, e, ph = t_all[valid], e_all[valid], p_all[valid]
    k = max(1, min(10, len(e)//5))
    sm = np.convolve(e, np.ones(k)/k, mode="same")
    avg = float(e.mean()); thr = avg - 15; lo = int(np.argmin(e))

    events, start = [], None
    for i, f in enumerate(sm < thr):
        if f and start is None:
            start = t[i]
        elif not f and start is not None:
            if t[i-1]-start >= 4:
                seg = e[(t >= start) & (t <= t[i-1])]
                events.append({"start": int(start), "end": int(t[i-1]),
                               "low": round(float(seg.min()), 1)})
            start = None
    if start is not None and t[-1]-start >= 4:
        events.append({"start": int(start), "end": int(t[-1]),
                       "low": round(float(e[t >= start].min()), 1)})

    zgrid = [[round(z["eng"]/z["total"]*100, 1) if z["total"] else None
              for z in row] for row in zones]
    return {"duration_s": int(duration), "average": round(avg, 1),
            "lowest_pct": round(float(e[lo]), 1), "lowest_time": int(t[lo]),
            "coverage": round(coverage, 1), "gap_seconds": gaps,
            "phone_peak": int(ph.max()) if len(ph) else 0,
            "phone_share": round(float((ph > 0).mean()*100), 1),
            "events": events, "zones": zgrid,
            "series": {"t": [int(x) for x in t], "e": [round(float(x), 1) for x in e],
                       "s": [round(float(x), 1) for x in sm]},
            "_log": {"t": [int(x) for x in t_all], "e": [round(float(x), 1) for x in e_all],
                     "p": [int(x) for x in p_all]}}


def new_zones():
    return [[{"eng": 0, "total": 0} for _ in range(GRID_COLS)] for _ in range(GRID_ROWS)]


def analyse_file(path, skip=None):
    cap = cv2.VideoCapture(path)
    if not cap.isOpened():
        return {"error": "Could not open that video file."}
    fps = cap.get(cv2.CAP_PROP_FPS) or 25
    src_h = cap.get(cv2.CAP_PROP_FRAME_HEIGHT) or 720
    src_w = cap.get(cv2.CAP_PROP_FRAME_WIDTH) or 1280
    if skip is None:
        skip = max(1, int(round(fps / SAMPLES_PER_SEC)))
    scale = min(1.0, INFER_WIDTH / max(src_w, 1))
    fw, fh = src_w * scale, src_h * scale
    total = int(cap.get(cv2.CAP_PROP_FRAME_COUNT)) or 0
    print(f"  analysing: {total} frames @ {fps:.0f}fps, "
          f"sampling every {skip} -> ~{total//max(skip,1)} inferences", flush=True)
    t_start = time.time()
    PROGRESS["stage"] = "Tracking students"
    hist = defaultdict(lambda: deque(maxlen=12)); ps = defaultdict(int)
    zones = new_zones()
    times, engs, phones = [], [], []
    last, idx, proc, cph = -1, 0, 0, []
    while True:
        ok, frame = cap.read()
        if not ok:
            break
        if idx % skip:
            idx += 1; continue
        if scale < 1.0:
            frame = cv2.resize(frame, (int(fw), int(fh)), interpolation=cv2.INTER_AREA)
        if proc % PHONE_EVERY == 0:
            cph = []
            for r in object_model(frame, classes=[PHONE_CLASS], conf=PHONE_CONF,
                                  verbose=False, imgsz=INFER_WIDTH):
                for b in r.boxes.xyxy:
                    x1, y1, x2, y2 = (float(v) for v in b)
                    cph.append(((x1+x2)/2, (y1+y2)/2))
        eng, dis, phn = judge_frame(frame, hist, ps, cph, fh, fw, zones)
        if proc % 10 == 0 and total:
            done = idx / total
            el = time.time() - t_start
            PROGRESS["percent"] = round(done * 100, 1)
            PROGRESS["elapsed"] = int(el)
            PROGRESS["eta"] = int(el / done - el) if done > 0.01 else 0
            if proc % 40 == 0:
                print(f"    {done*100:5.1f}%  ~{PROGRESS['eta']}s left", flush=True)
        tot = eng + dis
        sec = int(idx/fps)
        if sec != last:
            times.append(sec); engs.append((eng/tot*100) if tot else 0); phones.append(phn)
            last = sec
        proc += 1; idx += 1
    cap.release()
    PROGRESS.update({"percent": 100.0, "eta": 0, "stage": "Preparing the report"})
    if not times:
        return {"error": "No students were detected in this recording."}
    return summarise(times, engs, phones, zones, times[-1])


PROGRESS = {"active": False, "percent": 0.0, "eta": 0, "elapsed": 0,
            "stage": "", "filename": ""}


def reset_progress(filename=""):
    PROGRESS.update({"active": True, "percent": 0.0, "eta": 0, "elapsed": 0,
                     "stage": "Loading the recording", "filename": filename})


LIVE = {"on": False, "thread": None, "times": [], "engs": [], "phones": [],
        "zones": None, "started": 0, "err": None, "slides": []}


def live_loop(cam_index=0):
    cap = cv2.VideoCapture(cam_index)
    if not cap.isOpened():
        LIVE["err"] = "Could not open the camera."
        LIVE["on"] = False
        return
    src_h = cap.get(cv2.CAP_PROP_FRAME_HEIGHT) or 720
    src_w = cap.get(cv2.CAP_PROP_FRAME_WIDTH) or 1280
    scale = min(1.0, INFER_WIDTH / max(src_w, 1))
    fw, fh = src_w * scale, src_h * scale
    hist = defaultdict(lambda: deque(maxlen=12)); ps = defaultdict(int)
    LIVE["zones"] = new_zones()
    cph, proc, last = [], 0, -1
    next_sample = 0.0
    while LIVE["on"]:
        ok, frame = cap.read()
        if not ok:
            break
        now = time.time()
        if now < next_sample:
            continue
        next_sample = now + 1.0 / LIVE_SAMPLES_PS
        if scale < 1.0:
            frame = cv2.resize(frame, (int(fw), int(fh)), interpolation=cv2.INTER_AREA)
        if proc % PHONE_EVERY == 0:
            cph = []
            for r in object_model(frame, classes=[PHONE_CLASS], conf=PHONE_CONF,
                                  verbose=False, imgsz=INFER_WIDTH):
                for b in r.boxes.xyxy:
                    x1, y1, x2, y2 = (float(v) for v in b)
                    cph.append(((x1+x2)/2, (y1+y2)/2))
        eng, dis, phn = judge_frame(frame, hist, ps, cph, fh, fw, LIVE["zones"])
        tot = eng + dis
        sec = int(time.time() - LIVE["started"])
        if sec != last:
            LIVE["times"].append(sec)
            LIVE["engs"].append((eng/tot*100) if tot else 0)
            LIVE["phones"].append(phn)
            last = sec
        proc += 1
    cap.release()

@app.get("/")
def index():

    return FileResponse("Webapp/dashboard.html", headers={
        "Cache-Control": "no-store, no-cache, must-revalidate, max-age=0",
        "Pragma": "no-cache",
        "Expires": "0",
    })


@app.post("/api/signup")
def signup(resp: Response, name: str = Form(...), email: str = Form(...),
           password: str = Form(...)):
    missing = check_password(password)
    if missing:
        raise HTTPException(400, "Your password still needs " + ", ".join(missing) + ".")
    salt = secrets.token_hex(8)
    uid = uuid.uuid4().hex[:12]
    try:
        with db() as con:
            con.execute("INSERT INTO users VALUES(?,?,?,?,?,?)",
                        (uid, name.strip(), email.strip().lower(), salt,
                         hash_pw(password, salt), datetime.now().isoformat()))
    except sqlite3.IntegrityError:
        raise HTTPException(400, "An account already exists for that email.")
    tok = secrets.token_hex(16)
    with db() as con:
        con.execute("INSERT INTO tokens VALUES(?,?,?)",
                    (tok, uid, datetime.now().isoformat()))
    resp.set_cookie("lp_token", tok, httponly=True, samesite="lax")
    return {"name": name, "email": email}


@app.post("/api/login")
def login(resp: Response, email: str = Form(...), password: str = Form(...)):
    with db() as con:
        u = con.execute("SELECT * FROM users WHERE email=?",
                        (email.strip().lower(),)).fetchone()
    if not u or hash_pw(password, u["salt"]) != u["pwd"]:
        raise HTTPException(401, "That email and password do not match.")
    tok = secrets.token_hex(16)
    with db() as con:
        con.execute("INSERT INTO tokens VALUES(?,?,?)",
                    (tok, u["id"], datetime.now().isoformat()))
    resp.set_cookie("lp_token", tok, httponly=True, samesite="lax")
    return {"name": u["name"], "email": u["email"]}


@app.post("/api/logout")
def logout(resp: Response, lp_token: str = Cookie(None)):
    with db() as con:
        con.execute("DELETE FROM tokens WHERE token=?", (lp_token,))
    resp.delete_cookie("lp_token")
    return {"ok": True}


@app.get("/api/me")
def me(lp_token: str = Cookie(None)):
    u = user_from_token(lp_token)
    if not u:
        raise HTTPException(401, "Not signed in")
    return {"name": u["name"], "email": u["email"]}


@app.get("/api/sessions")
def list_sessions(lp_token: str = Cookie(None)):
    u = require(lp_token)
    with db() as con:
        rows = con.execute(
            "SELECT id,subject,class_name,mode,created,duration_s,average "
            "FROM sessions WHERE user_id=? ORDER BY created DESC", (u["id"],)).fetchall()
    return [dict(r) for r in rows]


@app.get("/api/sessions/{sid}")
def get_session(sid: str, lp_token: str = Cookie(None)):
    u = require(lp_token)
    with db() as con:
        r = con.execute("SELECT * FROM sessions WHERE id=? AND user_id=?",
                        (sid, u["id"])).fetchone()
    if not r:
        raise HTTPException(404, "Session not found")
    d = dict(r)
    d["events"] = json.loads(d.pop("events_json"))
    d["series"] = json.loads(d.pop("series_json"))
    d["zones"] = json.loads(d.pop("zones_json") or "[]")
    d["advice"] = advise(d)
    d["summary"] = build_summary(d)
    d["slides"] = json.loads(d.pop("slides_json") or "[]")
    d["slide_map"] = map_slides_to_engagement(d["slides"], d)
    d["content"] = analyse_slide_content(d["slide_map"])
    d["quiz"] = build_quiz(d)
    return d


def advise(d):
    """Threshold-based teaching guidance."""
    a, out = d["average"], []
    ev = d["events"]
    topics = [t for t in (d.get("topics") or "").split("|") if t.strip()]

    def topic_at(sec):

        best = None
        for t in topics:
            if ":" in t:
                m, nm = t.split(":", 1)
                try:
                    mm = float(m)
                except ValueError:
                    continue
                if mm * 60 <= sec and (best is None or mm * 60 > best[0]):
                    best = (mm * 60, nm.strip())
        return best[1] if best else None

    if a < 50:
        out.append({"level": "urgent", "title": "Attention was low for most of this class",
                    "body": f"Average engagement was {a}%, which indicates the session did not hold "
                            f"the room. This is usually structural rather than a single difficult "
                            f"topic: consider whether the class ran too long without interaction, "
                            f"whether the material was pitched above the group, or whether delivery "
                            f"was mostly one-directional."})
        if ev:
            w = min(ev, key=lambda e: e["low"])
            tp = topic_at(w["start"])
            out.append({"level": "urgent", "title": "The weakest stretch",
                        "body": f"Engagement bottomed out at {w['low']}% between "
                                f"{w['start']//60}m{w['start']%60:02d}s and "
                                f"{w['end']//60}m{w['end']%60:02d}s"
                                + (f", during \"{tp}\"." if tp else ".")
                                + " Re-teach this segment with a worked example and a check for "
                                  "understanding before moving on."})

    elif a < 70:

        out.append({"level": "warn", "title": "Attention held, but not firmly",
                    "body": f"Average engagement was {a}%. The class followed you, but not "
                            f"consistently, which most often indicates material that was "
                            f"understandable in principle but hard to picture. The suggestions "
                            f"below concern making the difficult passages more visual."})

        weak_topics = []
        for e in ev[:3]:
            tp = topic_at(e["start"])
            weak_topics.append((tp, e))

        if any(tp for tp, _ in weak_topics):
            for tp, e in weak_topics:
                if not tp:
                    continue
                out.append({"level": "warn", "title": f"Make \"{tp}\" more visual",
                            "body": f"Engagement fell to {e['low']}% while this topic was being "
                                    f"covered. Before teaching it again, consider replacing the "
                                    f"verbal explanation with something the class can see: a "
                                    f"diagram built up step by step on the board rather than "
                                    f"shown complete, a concrete worked example with real "
                                    f"numbers, or a simple analogy drawn from everyday "
                                    f"experience. Where a process is being described, a flow "
                                    f"diagram usually holds attention better than a list of "
                                    f"steps in prose."})
        elif ev:
            e = ev[0]
            out.append({"level": "warn", "title": "Make the difficult passage more visual",
                        "body": f"The clearest drop was between {e['start']//60}m"
                                f"{e['start']%60:02d}s and {e['end']//60}m{e['end']%60:02d}s, "
                                f"falling to {e['low']}%. Recall what you were explaining then "
                                f"and consider whether it could be shown rather than described: "
                                f"a diagram developed live on the board, a worked example with "
                                f"concrete values, a short demonstration, or an image standing "
                                f"in for an abstract idea. Abstract material delivered verbally "
                                f"is the most common cause of a drop of this shape."})

        out.append({"level": "warn", "title": "Break the exposition with something to look at",
                    "body": "As a general measure at this level of engagement, aim to put "
                            "something visual in front of the class every ten to fifteen "
                            "minutes - a diagram, a worked example, a short video, or a "
                            "question displayed for the room to consider. A change in what "
                            "the class is looking at restores attention more reliably than "
                            "a change in what they are hearing."})

    elif a < 80:
        out.append({"level": "good", "title": "The class held attention well",
                    "body": f"Average engagement was {a}%. Note what characterised this session "
                            f"- the pacing, the balance of activities, the material - so that it "
                            f"can be repeated deliberately rather than by chance."})

    else:
        out.append({"level": "good", "title": "A strong session",
                    "body": f"Average engagement was {a}%, which is high. The class stayed with "
                            f"you throughout, so the priority now is consolidation rather than "
                            f"correction: material delivered to an attentive room is best "
                            f"secured by a short recap while it is still fresh."})

        recap = []
        for t in topics:
            if ":" in t:
                m, nm = t.split(":", 1)
                recap.append(nm.strip())
        if recap:
            lines = "; ".join(f"({i+1}) {nm}" for i, nm in enumerate(recap[:6]))
            out.append({"level": "good", "title": "Suggested recap slides for the end of class",
                        "body": f"Based on the lecture plan you entered, a closing recap of five "
                                f"minutes could work through one slide per topic in the order "
                                f"taught: {lines}. Keep each slide to a single idea stated in "
                                f"one sentence with the supporting diagram beside it, and ask "
                                f"the class to supply the explanation before you reveal it. "
                                f"Reversing the direction in this way - showing the visual and "
                                f"asking for the concept, rather than stating the concept and "
                                f"showing the visual - tests recall rather than recognition."})
        else:
            out.append({"level": "good", "title": "Suggested recap for the end of class",
                        "body": "A closing recap of five minutes would consolidate a session "
                                "this strong. Prepare one slide for each main idea covered, each "
                                "carrying a single sentence and its supporting diagram, and work "
                                "through them in the order taught. Showing the diagram first and "
                                "asking the class to supply the idea, rather than the reverse, "
                                "tests recall rather than recognition. Entering a lecture plan "
                                "before the next class will allow this suggestion to name the "
                                "specific topics."})

        if ev:
            e = ev[0]
            out.append({"level": "info", "title": "One passage to include in the recap",
                        "body": f"Even in a strong session, engagement dipped to {e['low']}% "
                                f"between {e['start']//60}m{e['start']%60:02d}s and "
                                f"{e['end']//60}m{e['end']%60:02d}s"
                                + (f" during \"{topic_at(e['start'])}\"" if topic_at(e["start"]) else "")
                                + ". This is the passage most worth including in the recap."})

    for e in ev[:3]:
        tp = topic_at(e["start"])
        if tp:
            out.append({"level": "info", "title": f"Drop during \"{tp}\"",
                        "body": f"Engagement fell to {e['low']}% between "
                                f"{e['start']//60}m{e['start']%60:02d}s and "
                                f"{e['end']//60}m{e['end']%60:02d}s. Worth revisiting."})

    if d.get("phone_share", 0) > 15:
        out.append({"level": "info", "title": "Phone use was noticeable",
                    "body": f"Phones were visible for {d['phone_share']}% of the session, "
                            f"peaking at {d['phone_peak']} at once. Where this overlaps the "
                            f"drops above it is more likely a symptom than a cause; a "
                            f"signposted break or a device-based activity works better than "
                            f"a prohibition."})

    z = d.get("zones") or []
    flat = [(r, c, v) for r, row in enumerate(z) for c, v in enumerate(row) if v is not None]
    if len(flat) >= 4:
        worst = min(flat, key=lambda x: x[2])
        best = max(flat, key=lambda x: x[2])
        if best[2] - worst[2] > 20:
            rows = ["front", "middle", "back"]
            cols = ["far left", "left", "right", "far right"]
            out.append({"level": "info", "title": "One part of the room lagged",
                        "body": f"The {rows[min(worst[0],2)]} {cols[min(worst[1],3)]} zone "
                                f"averaged {worst[2]}% against {best[2]}% in the strongest "
                                f"zone. Moving into that part of the room, or directing "
                                f"questions there, usually closes the gap."})
    return out


def parse_slides(path, filename):
    """Extract a title and a short text sample from each slide or page.

    Text is EXTRACTED from the file, never generated: the system reports what the
    slides say, not what it supposes they mean.
    """
    ext = os.path.splitext(filename or "")[1].lower()
    slides = []

    if ext == ".pptx":
        try:
            from pptx import Presentation
        except ImportError:
            return [], "python-pptx is not installed. Run: pip install python-pptx"
        try:
            prs = Presentation(path)
        except Exception as e:
            return [], f"Could not read that PowerPoint file: {e}"
        for i, slide in enumerate(prs.slides, 1):
            title, body = "", []
            try:
                if slide.shapes.title is not None and slide.shapes.title.text.strip():
                    title = slide.shapes.title.text.strip()
            except Exception:
                pass
            for sh in slide.shapes:
                if sh.has_text_frame:
                    for para in sh.text_frame.paragraphs:
                        t = "".join(r.text for r in para.runs).strip()
                        if t and t != title:
                            body.append(t)
            if not title:
                title = body[0] if body else f"Slide {i}"
                body = body[1:]
            slides.append({"n": i, "title": title[:120],
                           "text": " · ".join(body)[:300]})

    elif ext == ".pdf":
        try:
            from pypdf import PdfReader
        except ImportError:
            return [], "pypdf is not installed. Run: pip install pypdf"
        try:
            reader = PdfReader(path)
        except Exception as e:
            return [], f"Could not read that PDF: {e}"
        for i, page in enumerate(reader.pages, 1):
            try:
                raw = page.extract_text() or ""
            except Exception:
                raw = ""
            lines = [l.strip() for l in raw.split("\n") if l.strip()]
            title = lines[0][:120] if lines else f"Page {i}"
            slides.append({"n": i, "title": title,
                           "text": " · ".join(lines[1:])[:300]})
    else:
        return [], "Please upload a .pptx or .pdf file."

    if not slides:
        return [], "No slides or pages could be read from that file."
    return slides, None


def map_slides_to_engagement(slides, d):
    """Attach an engagement figure to each slide.

    ASSUMPTION, stated to the user: slides are assumed to have been shown in order
    and for equal durations across the session. The system has no way to observe
    when a slide was actually displayed, so this is an approximation, not a
    measurement, and it is labelled as such in the interface.
    """
    if not slides:
        return []
    ts = (d.get("series") or {}).get("t") or []
    es = (d.get("series") or {}).get("e") or []
    total = d.get("duration_s") or (ts[-1] if ts else 0)
    if not ts or not total:
        return [dict(sl, start=None, end=None, engagement=None, drops=0) for sl in slides]

    ev = d.get("events") or []
    per = total / len(slides)
    out = []
    for i, sl in enumerate(slides):
        start, end = i * per, (i + 1) * per
        vals = [e for t, e in zip(ts, es) if start <= t < end]
        avg = round(float(np.mean(vals)), 1) if vals else None
        drops = sum(1 for x in ev if start <= x["start"] < end)
        out.append(dict(sl, start=int(start), end=int(end),
                        engagement=avg, drops=drops))
    return out

OLLAMA_URL = os.environ.get("OLLAMA_URL", "http://localhost:11434")
OLLAMA_MODEL = os.environ.get("OLLAMA_MODEL", "llama3.2")
OPENAI_KEY = os.environ.get("OPENAI_API_KEY")
OPENAI_BASE = os.environ.get("OPENAI_BASE", "https://api.openai.com/v1")
OPENAI_MODEL = os.environ.get("OPENAI_MODEL", "gpt-4o-mini")


def llm_status():
    """Report which provider, if any, is available."""
    import urllib.request
    try:
        with urllib.request.urlopen(f"{OLLAMA_URL}/api/tags", timeout=2) as r:
            tags = json.loads(r.read().decode())
        names = [m.get("name", "") for m in tags.get("models", [])]
        if names:
            chosen = next((n for n in names if n.startswith(OLLAMA_MODEL)), names[0])
            return {"available": True, "provider": "ollama", "model": chosen,
                    "local": True}
    except Exception:
        pass
    if OPENAI_KEY:
        return {"available": True, "provider": "openai", "model": OPENAI_MODEL,
                "local": False}
    return {"available": False, "provider": None, "model": None, "local": None,
            "note": ("No language model is configured. For a local model that keeps "
                     "lecture content on this machine, install Ollama from ollama.com "
                     "and run: ollama pull llama3.2. Alternatively set OPENAI_API_KEY "
                     "to use a hosted model, bearing in mind that slide text would then "
                     "be sent to a third party.")}


def llm_complete(prompt, timeout=180):
    """Send one prompt to whichever provider is available. Returns text or None."""
    import urllib.request
    st = llm_status()
    if not st["available"]:
        return None, st.get("note")

    if st["provider"] == "ollama":
        body = json.dumps({"model": st["model"], "prompt": prompt, "stream": False,
                           "options": {"temperature": 0.2}}).encode()
        req = urllib.request.Request(f"{OLLAMA_URL}/api/generate", data=body,
                                     headers={"Content-Type": "application/json"})
        try:
            with urllib.request.urlopen(req, timeout=timeout) as r:
                return json.loads(r.read().decode()).get("response", "").strip(), None
        except Exception as e:
            return None, f"The local model did not respond: {e}"

    body = json.dumps({"model": st["model"], "temperature": 0.2,
                       "messages": [{"role": "user", "content": prompt}]}).encode()
    req = urllib.request.Request(f"{OPENAI_BASE}/chat/completions", data=body,
                                 headers={"Content-Type": "application/json",
                                          "Authorization": f"Bearer {OPENAI_KEY}"})
    try:
        with urllib.request.urlopen(req, timeout=timeout) as r:
            data = json.loads(r.read().decode())
        return data["choices"][0]["message"]["content"].strip(), None
    except Exception as e:
        return None, f"The hosted model did not respond: {e}"


def llm_class_recap(slide_map, d):
    """Produce a content recap of the class from the slides and the engagement data."""
    if not slide_map:
        return {"ok": False, "error": "No slides were uploaded for this session."}

    lines = []
    for sl in slide_map[:60]:
        eng = sl.get("engagement")
        lines.append(f"[Slide {sl['n']}] {sl.get('title','')}"
                     + (f" (class engagement {eng}%)" if eng is not None else "")
                     + (f"\n    {sl.get('text','')[:220]}" if sl.get("text") else ""))
    deck = "\n".join(lines)

    engs = [sl["engagement"] for sl in slide_map if sl.get("engagement") is not None]
    avg = round(sum(engs) / len(engs), 1) if engs else None
    weak = [sl for sl in slide_map if sl.get("engagement") is not None
            and avg is not None and sl["engagement"] < avg - 10]
    weak_txt = ", ".join(f"slide {sl['n']} ({sl.get('title','')}, {sl['engagement']}%)"
                         for sl in weak[:8]) or "none"

    prompt = f"""You are helping a university lecturer review a class they have just taught.

Below are the slides they used. Where a figure is given, it is the measured class engagement while that slide was on screen (higher is better). The session averaged {avg}%.

Slides where engagement fell more than ten points below average: {weak_txt}

SLIDES:
{deck}

Write a recap for the lecturer with exactly these four sections, using these exact headings:

SUMMARY
Three or four sentences describing what this class covered, based only on the slides. Do not speculate about material not present in them.

CONCEPTS
The five or six main concepts taught, one per line, each formatted as "Concept name — one short sentence explaining it in plain language."

WEAK POINTS
For each slide listed above as low engagement, one line formatted as "Slide N, Concept — why this material is typically difficult, and one concrete suggestion for presenting it more clearly." If no slides were listed, write "None identified."

RECAP PLAN
Four or five bullet points a lecturer could use as a five-minute closing recap of this class, in the order taught.

Write plainly and concretely. Do not invent content that is not in the slides."""

    text, err = llm_complete(prompt)
    if text is None:
        return {"ok": False, "error": err or "The model returned nothing."}

    sections = split_sections(text, ["SUMMARY", "CONCEPTS", "WEAK POINTS", "RECAP PLAN"])
    st = llm_status()
    parsed_anything = any(sections.get(k) for k in ["SUMMARY", "CONCEPTS", "WEAK POINTS", "RECAP PLAN"])
    return {"ok": True, "model": st.get("model"), "local": st.get("local"),
            "summary": " ".join(sections.get("SUMMARY", [])),
            "concepts": sections.get("CONCEPTS", []),
            "weak": sections.get("WEAK POINTS", []),
            "plan": sections.get("RECAP PLAN", []),
            "raw": None if parsed_anything else text}


def analyse_slide_content(slide_map):
    """Statistical content analysis of the uploaded slides.

    Key terms are EXTRACTED from the slide text by term-frequency weighting; the
    system does not read, interpret or summarise the material. The value added
    here is not comprehension but correlation: terms are cross-referenced against
    the engagement recorded while their slide was on screen, which identifies the
    concepts that were being presented when the class lost attention.
    """
    if not slide_map:
        return None
    docs, idx = [], []
    for sl in slide_map:
        text = f"{sl.get('title','')} {sl.get('text','')}".strip()
        if len(text) > 3:
            docs.append(text)
            idx.append(sl)
    if len(docs) < 2:
        return None

    try:
        from sklearn.feature_extraction.text import TfidfVectorizer
    except ImportError:
        return None

    STOP = ("english")
    try:
        vec = TfidfVectorizer(stop_words=STOP, ngram_range=(1, 2),
                              max_features=400, token_pattern=r"[A-Za-z][A-Za-z\-]{2,}")
        M = vec.fit_transform(docs)
    except ValueError:
        return None
    terms = np.array(vec.get_feature_names_out())
    if len(terms) == 0:
        return None

    overall = np.asarray(M.sum(axis=0)).ravel()
    order = np.argsort(-overall)[:18]
    key_terms = [{"term": terms[i], "weight": round(float(overall[i]), 3)}
                 for i in order if overall[i] > 0]

    per_slide = []
    for r, sl in enumerate(idx):
        row = np.asarray(M[r].todense()).ravel()
        top = np.argsort(-row)[:4]
        per_slide.append({
            "n": sl["n"], "title": sl.get("title", ""),
            "engagement": sl.get("engagement"),
            "terms": [terms[i] for i in top if row[i] > 0],
        })

    engs = [sl.get("engagement") for sl in idx if sl.get("engagement") is not None]
    at_risk = []
    if engs:
        cutoff = float(np.mean(engs)) - 10
        weak_rows = [r for r, sl in enumerate(idx)
                     if sl.get("engagement") is not None and sl["engagement"] < cutoff]
        if weak_rows:
            weak = np.asarray(M[weak_rows].sum(axis=0)).ravel()
            strong_rows = [r for r in range(len(idx)) if r not in weak_rows]
            strong = (np.asarray(M[strong_rows].sum(axis=0)).ravel()
                      if strong_rows else np.zeros_like(weak))
            diff = weak - strong * (len(weak_rows) / max(len(strong_rows), 1))
            for i in np.argsort(-diff)[:8]:
                if weak[i] > 0 and diff[i] > 0:
                    slides_with = [idx[r]["n"] for r in weak_rows
                                   if np.asarray(M[r].todense()).ravel()[i] > 0]
                    at_risk.append({"term": terms[i], "slides": slides_with[:5]})

    return {"key_terms": key_terms, "per_slide": per_slide, "at_risk": at_risk,
            "cutoff_note": ("Concepts below are those appearing distinctively on slides "
                            "where engagement fell more than ten points below this "
                            "session's average.")}


def build_slide_quiz(slides, content):
    """Create an answerable quiz using only information extracted from the deck."""
    if not slides:
        return {"available": False, "items": [],
                "note": "No readable slides were available for a quiz."}

    per_slide = {p["n"]: p.get("terms") or []
                 for p in (content.get("per_slide") or [])}
    all_terms = []
    for p in content.get("per_slide") or []:
        for term in p.get("terms") or []:
            if term and term not in all_terms:
                all_terms.append(term)

    items = []
    for sl in slides:
        terms = per_slide.get(sl["n"], [])
        if not terms:
            continue
        correct = terms[0]
        distractors = [t for t in all_terms if t != correct and t not in terms[:2]][:3]
        if len(distractors) < 3:
            continue
        options = [correct] + distractors
        shift = int(sl["n"]) % len(options)
        options = options[shift:] + options[:shift]
        items.append({
            "question": f"Which key concept is most strongly associated with slide "
                        f"{sl['n']}, ‘{sl.get('title') or 'Untitled'}’?",
            "options": options,
            "answer": options.index(correct),
            "explanation": f"‘{correct}’ is the most prominent extracted term on "
                           f"slide {sl['n']}.",
        })
        if len(items) == 6:
            break

    titles = [(sl["n"], (sl.get("title") or "").strip()) for sl in slides]
    titles = [(n, t) for n, t in titles if t and t.lower() != f"slide {n}"]
    unique_titles = []
    for _, title in titles:
        if title not in unique_titles:
            unique_titles.append(title)
    if len(items) < 3 and len(unique_titles) >= 4:
        items = []
        for n, correct in titles[:6]:
            distractors = [t for t in unique_titles if t != correct][:3]
            if len(distractors) < 3:
                continue
            options = [correct] + distractors
            shift = int(n) % len(options)
            options = options[shift:] + options[:shift]
            items.append({
                "question": f"Which title belongs to slide {n}?",
                "options": options,
                "answer": options.index(correct),
                "explanation": f"Slide {n} is titled ‘{correct}’.",
            })

    return {
        "available": bool(items),
        "items": items,
        "note": ("Questions are created directly from the uploaded slide titles and "
                 "statistically extracted terms. The slide file is deleted after parsing."
                 if items else
                 "This deck did not contain enough distinct readable text to create a "
                 "reliable multiple-choice quiz."),
    }


def split_sections(text, headings):
    """Split model output on headings, tolerating markdown and numbering.

    Small models rarely reproduce a heading exactly as asked: they wrap it in
    asterisks, prefix it with a number, or append a colon. Normalising each line
    before comparison makes the parser robust to all of these variations.
    """
    import re

    def norm(line):
        t = re.sub(r"^[#>\s]*", "", line.strip())
        t = re.sub(r"^\d+[\.\)]\s*", "", t)
        return t.strip("*_ ").rstrip(":").strip().upper()

    wanted = {h.upper() for h in headings}
    out, current = {}, None
    for raw in text.split("\n"):
        key = norm(raw)
        if key in wanted:
            current = key
            out[current] = []
            continue
        line = raw.strip()
        if current and line:
            out[current].append(re.sub(r"^[-•*\d\.\)\s]+", "", line).strip())
    return out


def llm_slides_recap(slides):
    if not slides:
        return {"ok": False, "error": "No slides were read from that file."}
    lines = []
    for sl in slides[:60]:
        lines.append(f"[{sl['n']}] {sl.get('title','')}"
                     + (f"\n    {sl.get('text','')[:250]}" if sl.get("text") else ""))
    deck = "\n".join(lines)

    prompt = f"""You are helping a university lecturer review the material they taught.

Below are their lecture slides, given as slide number, title and the text on the slide.

SLIDES:
{deck}

Write a recap using exactly these three headings:

SUMMARY
Four or five sentences describing what this material covers and how it is structured, based only on the slides. Do not speculate about anything not present in them.

CONCEPTS
The six to eight main concepts, one per line, each formatted as "Concept name — one clear sentence explaining it in plain language a student could follow."

RECAP PLAN
Five or six bullet points forming a five-minute closing recap of this material, in the order it was taught, each naming what to show and what to ask the class.

Write plainly and concretely. Do not invent content that is not in the slides."""

    text, err = llm_complete(prompt)
    if text is None:
        return {"ok": False, "error": err or "The model returned nothing."}

    sections = split_sections(text, ["SUMMARY", "CONCEPTS", "RECAP PLAN"])
    st = llm_status()
    parsed_anything = any(sections.get(k) for k in ["SUMMARY", "CONCEPTS", "RECAP PLAN"])
    return {"ok": True, "model": st.get("model"), "local": st.get("local"),
            "summary": " ".join(sections.get("SUMMARY", [])),
            "concepts": sections.get("CONCEPTS", []),
            "plan": sections.get("RECAP PLAN", []),
            "raw": None if parsed_anything else text}


def build_summary(d):
    a = d["average"]
    ev = d.get("events") or []
    zones = d.get("zones") or []
    series = (d.get("series") or {}).get("s") or (d.get("series") or {}).get("e") or []
    topics = [t for t in (d.get("topics") or "").split("|") if t.strip()]

    def topic_at(sec):
        best = None
        for t in topics:
            if ":" in t:
                m, nm = t.split(":", 1)
                try:
                    mm = float(m)
                except ValueError:
                    continue
                if mm * 60 <= sec and (best is None or mm * 60 > best[0]):
                    best = (mm * 60, nm.strip())
        return best[1] if best else None

    def clk(sec):
        sec = int(sec)
        return f"{sec//60}m {sec%60:02d}s"

    if a >= 80:
        verdict = ("The class stayed with you", 
                   f"Average engagement was {a}%. This was a strong session.")
    elif a >= 70:
        verdict = ("The class largely stayed with you",
                   f"Average engagement was {a}%, which is solid.")
    elif a >= 50:
        verdict = ("Attention held, but not firmly",
                   f"Average engagement was {a}%. The room followed you, but inconsistently.")
    else:
        verdict = ("Attention was low for much of the class",
                   f"Average engagement was {a}%. The session did not hold the room.")

    shape = None
    if len(series) >= 12:
        n = len(series)
        first, mid, last = (float(np.mean(series[:n//3])),
                            float(np.mean(series[n//3:2*n//3])),
                            float(np.mean(series[2*n//3:])))
        if last < first - 8 and last < mid - 4:
            shape = ("Attention declined toward the end",
                     f"The class averaged {first:.0f}% in the opening third and "
                     f"{last:.0f}% in the closing third. This is the classic pattern of "
                     f"attention decay in a long session, and it usually responds to a "
                     f"break or a change of activity rather than a change of content.")
        elif first < last - 8:
            shape = ("Attention built through the session",
                     f"The class began at {first:.0f}% and finished at {last:.0f}%. "
                     f"Whatever changed after the opening is worth repeating earlier "
                     f"next time.")
        elif mid < first - 8 and mid < last - 8:
            shape = ("Attention dipped in the middle",
                     f"The opening ({first:.0f}%) and close ({last:.0f}%) held better "
                     f"than the middle ({mid:.0f}%). The material covered mid-session is "
                     f"the place to look.")
        else:
            shape = ("Attention was steady throughout",
                     f"The class stayed close to its average across all three thirds "
                     f"({first:.0f}%, {mid:.0f}%, {last:.0f}%), with no systematic drift.")

    weakest = None
    if ev:
        w = min(ev, key=lambda e: e["low"])
        tp = topic_at(w["start"])
        weakest = ("Weakest stretch",
                   f"Engagement fell to {w['low']}% between {clk(w['start'])} and "
                   f"{clk(w['end'])}"
                   + (f", during \"{tp}\"" if tp else "")
                   + f". Of the {len(ev)} flagged period"
                   + ("s" if len(ev) != 1 else "")
                   + ", this is the one to revisit first.")

    strongest = None
    if series:
        peak_i = int(np.argmax(series))
        times = (d.get("series") or {}).get("t") or []
        if times and peak_i < len(times):
            pt = times[peak_i]
            tp = topic_at(pt)
            strongest = ("Strongest moment",
                         f"Engagement peaked at {series[peak_i]:.0f}% around {clk(pt)}"
                         + (f", during \"{tp}\"" if tp else "")
                         + ". Whatever held the room here is worth using more often.")

    spatial = None
    flat = [(r, c, v) for r, row in enumerate(zones) for c, v in enumerate(row)
            if v is not None]
    if len(flat) >= 4:
        worst = min(flat, key=lambda x: x[2]); best = max(flat, key=lambda x: x[2])
        rows = ["front", "middle", "back"]; cols = ["far left", "left", "right", "far right"]
        if best[2] - worst[2] > 20:
            spatial = ("One part of the room lagged",
                       f"The {rows[min(worst[0],2)]} {cols[min(worst[1],3)]} area averaged "
                       f"{worst[2]}% against {best[2]}% in the strongest zone. Moving into "
                       f"that part of the room, or directing questions there, usually "
                       f"closes the gap.")
        else:
            spatial = ("Attention was evenly spread",
                       f"No part of the room lagged materially: zones ranged from "
                       f"{worst[2]}% to {best[2]}%.")

    phones = None
    if d.get("phone_share", 0) > 10:
        phones = ("Phones were visible",
                  f"Devices were in view for {d['phone_share']}% of the session, peaking at "
                  f"{d['phone_peak']} at once. Where this overlaps the drops above it is "
                  f"more likely a symptom than a cause.")

    if a < 50:
        carry = ("Take this into the next class",
                 "Reconsider the structure of the session before reconsidering the content: "
                 "shorter exposition, more frequent interaction.")
    elif a < 70 and ev:
        w = min(ev, key=lambda e: e["low"]); tp = topic_at(w["start"])
        carry = ("Take this into the next class",
                 f"Open with a short revisit of "
                 + (f"\"{tp}\"" if tp else f"the material from around {clk(w['start'])}")
                 + ", using something the class can see rather than a verbal explanation.")
    elif ev:
        w = min(ev, key=lambda e: e["low"]); tp = topic_at(w["start"])
        carry = ("Take this into the next class",
                 f"The session worked overall; the one passage worth a brief recap is "
                 + (f"\"{tp}\"" if tp else f"the material around {clk(w['start'])}") + ".")
    else:
        carry = ("Take this into the next class",
                 "Nothing needs correcting. Note what characterised this session so it can "
                 "be repeated deliberately.")

    cards = [verdict]
    for c in (shape, weakest, strongest, spatial, phones, carry):
        if c:
            cards.append(c)
    return [{"title": t, "body": b} for t, b in cards]


def build_quiz(d):
    topics, source = [], "plan"
    for t in (d.get("topics") or "").split("|"):
        if ":" in t:
            m, nm = t.split(":", 1)
            try:
                topics.append((float(m) * 60, nm.strip()))
            except ValueError:
                continue

    if not topics:
        source = "slides"
        smap = d.get("slide_map") or []
        if smap:
            step = max(1, len(smap) // 8)
            for i in range(0, len(smap), step):
                block = smap[i:i + step]
                title = (block[0].get("title") or f"Slides {block[0]['n']}").strip()
                start = block[0].get("start")
                topics.append((float(start if start is not None else i), title))

    if not topics:
        return {"available": False,
                "note": "No topics could be identified. Enter a lecture plan, or upload "
                        "slides with readable titles, and a recap quiz will be prepared "
                        "here, ordered by the material that lost the most attention."}

    topics.sort()
    ev = d.get("events") or []
    series = (d.get("series") or {}); ts = series.get("t") or []; es = series.get("e") or []

    scored = []
    for i, (start, name) in enumerate(topics):
        end = topics[i + 1][0] if i + 1 < len(topics) else (ts[-1] if ts else start + 600)
        vals = [e for t, e in zip(ts, es) if start <= t < end]
        avg = float(np.mean(vals)) if vals else None
        drops = sum(1 for x in ev if start <= x["start"] < end)
        scored.append({"topic": name, "start": int(start), "end": int(end),
                       "engagement": round(avg, 1) if avg is not None else None,
                       "drops": drops})

    ranked = sorted(scored, key=lambda x: (x["engagement"] if x["engagement"] is not None else 999))

    STEMS = [
        "In one sentence, what is {t} and why does it matter?",
        "Give one concrete example of {t} and explain what makes it an example.",
        "What is the most common mistake people make with {t}?",
        "How does {t} relate to what we covered before it?",
        "If you had to explain {t} to someone who missed the class, what would you say first?",
    ]

    items = []
    for rank, sc in enumerate(ranked):
        if sc["engagement"] is None:
            priority = "unknown"
        elif sc["engagement"] < 50 or sc["drops"] >= 2:
            priority = "high"
        elif sc["engagement"] < 70 or sc["drops"] == 1:
            priority = "medium"
        else:
            priority = "low"
        items.append({
            "topic": sc["topic"],
            "engagement": sc["engagement"],
            "drops": sc["drops"],
            "priority": priority,
            "window": f"{sc['start']//60}m–{sc['end']//60}m",
            "stems": [STEMS[rank % len(STEMS)],
                      STEMS[(rank + 2) % len(STEMS)]],
        })

    src_note = ("your lecture plan" if source == "plan"
                else "the titles of your uploaded slides")
    return {"available": True, "items": items, "source": source,
            "note": f"Topics are taken from {src_note}. These are question stems, not "
                    "finished questions: the system knows the topic names but not what you "
                    "taught, so the subject content is yours to supply. The ordering is "
                    "data-driven — the material with the weakest engagement appears first, "
                    "so what most needs checking is at the top."}


def save_session(user, subject, class_name, topics, mode, res):
    sid = uuid.uuid4().hex[:12]
    log = res.pop("_log", None)
    if log:
        with open(os.path.join(LOG_DIR, f"{sid}.csv"), "w", newline="") as f:
            f.write("time_seconds,engaged_percent,phones_visible\n")
            for a, b, c in zip(log["t"], log["e"], log["p"]):
                f.write(f"{a},{b},{c}\n")
    with db() as con:
        con.execute("""INSERT INTO sessions
            (id,user_id,subject,class_name,topics,mode,created,duration_s,average,
             lowest_pct,lowest_time,coverage,gap_seconds,phone_peak,phone_share,
             events_json,series_json,zones_json,slides_json)
            VALUES(?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?)""",
                    (sid, user["id"], subject, class_name, topics, mode,
                     datetime.now().isoformat(timespec="minutes"),
                     res["duration_s"], res["average"], res["lowest_pct"],
                     res["lowest_time"], res["coverage"], res["gap_seconds"],
                     res["phone_peak"], res["phone_share"],
                     json.dumps(res["events"]), json.dumps(res["series"]),
                     json.dumps(res["zones"]), json.dumps(res.get("slides") or [])))
    res["id"] = sid
    res["subject"] = subject
    res["class_name"] = class_name
    res["topics"] = topics
    res["advice"] = advise(res)
    res["summary"] = build_summary(res)
    res["slide_map"] = map_slides_to_engagement(res.get("slides") or [], res)
    res["content"] = analyse_slide_content(res["slide_map"])
    res["quiz"] = build_quiz(res)
    return res


@app.post("/api/analyze")
async def api_analyze(file: UploadFile = File(...), subject: str = Form("Class"),
                      class_name: str = Form(""), topics: str = Form(""),
                      slides: UploadFile = File(None),
                      lp_token: str = Cookie(None)):
    u = require(lp_token)
    if slides is None or not slides.filename:
        raise HTTPException(400, "Lecture slides are required. Please attach the "
                                 "PDF or PowerPoint used in this class before analysing.")
    parsed_slides, slide_error = [], None
    if slides is not None and slides.filename:
        sp = os.path.join(UPLOAD_DIR, f"{uuid.uuid4().hex}_{slides.filename}")
        with open(sp, "wb") as out:
            shutil.copyfileobj(slides.file, out)
        try:
            parsed_slides, slide_error = parse_slides(sp, slides.filename)
        finally:
            if os.path.exists(sp):
                os.remove(sp)   
    ext = os.path.splitext(file.filename or "v.mp4")[1] or ".mp4"
    tmp = os.path.join(UPLOAD_DIR, f"{uuid.uuid4().hex}{ext}")
    with open(tmp, "wb") as out:
        shutil.copyfileobj(file.file, out)
    reset_progress(file.filename or "recording")
    try:
        res = analyse_file(tmp)
    finally:
        PROGRESS["active"] = False
        if os.path.exists(tmp):
            os.remove(tmp)       
    if "error" in res:
        return JSONResponse(res, status_code=400)
    if slide_error:
        raise HTTPException(400, slide_error)
    if not parsed_slides:
        raise HTTPException(400, "No slides could be read from that file. Please check "
                                 "it is a valid .pptx or .pdf.")
    res["slides"] = parsed_slides
    out = save_session(u, subject, class_name or (file.filename or "Session"),
                       topics, "recorded", res)
    if slide_error:
        out["slide_error"] = slide_error
    return out


@app.get("/api/progress")
def progress(lp_token: str = Cookie(None)):
    require(lp_token)
    return dict(PROGRESS)


@app.post("/api/slides/parse")
async def api_slides_parse(slides: UploadFile = File(...), lp_token: str = Cookie(None)):
    """Read a slide deck and describe its content. No video, no timing."""
    require(lp_token)
    sp = os.path.join(UPLOAD_DIR, f"{uuid.uuid4().hex}_{slides.filename}")
    with open(sp, "wb") as out:
        shutil.copyfileobj(slides.file, out)
    try:
        parsed, err = parse_slides(sp, slides.filename)
    finally:
        if os.path.exists(sp):
            os.remove(sp)          # parsed, then deleted
    if err:
        raise HTTPException(400, err)
    plain = [dict(sl, engagement=None, drops=0, start=None, end=None) for sl in parsed]
    content = analyse_slide_content(plain)
    return {"filename": slides.filename, "slides": parsed,
            "content": content, "quiz": build_slide_quiz(parsed, content)}


@app.post("/api/slides/ai")
async def api_slides_ai(payload: dict, lp_token: str = Cookie(None)):
    require(lp_token)
    return llm_slides_recap(payload.get("slides") or [])


@app.get("/api/llm/status")
def api_llm_status(lp_token: str = Cookie(None)):
    require(lp_token)
    return llm_status()


@app.post("/api/sessions/{sid}/recap")
def api_recap(sid: str, lp_token: str = Cookie(None)):
    u = require(lp_token)
    with db() as con:
        r = con.execute("SELECT * FROM sessions WHERE id=? AND user_id=?",
                        (sid, u["id"])).fetchone()
    if not r:
        raise HTTPException(404, "Session not found")
    d = dict(r)
    d["events"] = json.loads(d.pop("events_json") or "[]")
    d["series"] = json.loads(d.pop("series_json") or "{}")
    d["zones"] = json.loads(d.pop("zones_json") or "[]")
    slides = json.loads(d.pop("slides_json") or "[]")
    smap = map_slides_to_engagement(slides, d)
    return llm_class_recap(smap, d)


@app.post("/api/live/start")
async def live_start(slides: UploadFile = File(...), lp_token: str = Cookie(None)):
    require(lp_token)
    if LIVE["on"]:
        raise HTTPException(400, "A live class is already running.")
    if slides is None or not slides.filename:
        raise HTTPException(400, "Lecture slides are required before starting a class.")
    sp = os.path.join(UPLOAD_DIR, f"{uuid.uuid4().hex}_{slides.filename}")
    with open(sp, "wb") as out:
        shutil.copyfileobj(slides.file, out)
    try:
        parsed, err = parse_slides(sp, slides.filename)
    finally:
        if os.path.exists(sp):
            os.remove(sp)
    if err:
        raise HTTPException(400, err)
    if not parsed:
        raise HTTPException(400, "No slides could be read from that file.")
    LIVE.update({"on": True, "times": [], "engs": [], "phones": [],
                 "zones": new_zones(), "started": time.time(), "err": None,
                 "slides": parsed})
    LIVE["thread"] = threading.Thread(target=live_loop, daemon=True)
    LIVE["thread"].start()
    time.sleep(1.5)
    if LIVE["err"]:
        LIVE["on"] = False
        raise HTTPException(400, LIVE["err"])
    return {"started": True}


@app.get("/api/live/status")
def live_status(lp_token: str = Cookie(None)):
    require(lp_token)
    e = LIVE["engs"]
    recent = e[-60:]
    return {"on": LIVE["on"], "error": LIVE["err"],
            "elapsed": int(time.time() - LIVE["started"]) if LIVE["on"] else 0,
            "current": round(float(recent[-1]), 1) if recent else 0,
            "average": round(float(np.mean(e)), 1) if e else 0,
            "phones": LIVE["phones"][-1] if LIVE["phones"] else 0,
            "series": {"t": LIVE["times"][-180:],
                       "e": [round(float(v), 1) for v in e[-180:]]}}


@app.post("/api/live/stop")
def live_stop(subject: str = Form("Class"), class_name: str = Form(""),
              topics: str = Form(""), lp_token: str = Cookie(None)):
    u = require(lp_token)
    if not LIVE["on"]:
        raise HTTPException(400, "No live class is running.")
    LIVE["on"] = False
    time.sleep(1.0)
    if not LIVE["times"]:
        raise HTTPException(400, "No students were detected during the class.")
    res = summarise(LIVE["times"], LIVE["engs"], LIVE["phones"],
                    LIVE["zones"], LIVE["times"][-1])
    if "error" in res:
        return JSONResponse(res, status_code=400)
    res["slides"] = LIVE.get("slides") or []
    return save_session(u, subject, class_name or "Live class", topics, "live", res)


def build_pdf(sid, label):
    """Generate the PDF in-process (no subprocess) and return its path."""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import mm
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer, Table,
                                    TableStyle, Image)
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    with db() as con:
        r = con.execute("SELECT * FROM sessions WHERE id=?", (sid,)).fetchone()
    if not r:
        raise HTTPException(404, "Session not found.")
    d = dict(r)
    d["events"] = json.loads(d.pop("events_json") or "[]")
    d["series"] = json.loads(d.pop("series_json") or "{}")
    d["zones"] = json.loads(d.pop("zones_json") or "[]")
    advice = advise(d)

    safe = "".join(ch for ch in label if ch.isalnum() or ch in " -_").strip() or "session"
    out_dir = tempfile.mkdtemp()
    pdf_path = os.path.join(out_dir, f"LearnPulse_Report_{safe}.pdf")
    chart_path = os.path.join(out_dir, "chart.png")

    def fmt(sec):
        sec = int(sec)
        return f"{sec//60}m {sec%60:02d}s"

    t = d["series"].get("t", [])
    e = d["series"].get("e", [])
    sm = d["series"].get("s", e)
    fig, ax = plt.subplots(figsize=(7.2, 2.8))
    if t:
        mins = [x / 60 for x in t]
        ax.plot(mins, e, color="#C9C2D6", linewidth=1)
        ax.plot(mins, sm, color="#5B3E96", linewidth=2.2)
        ax.axhline(d["average"], color="#999999", linestyle="--", linewidth=1)
        for ev in d["events"]:
            ax.axvspan(ev["start"] / 60, ev["end"] / 60, color="#B8860B", alpha=.15)
    ax.set_xlabel("Time (minutes)"); ax.set_ylabel("Class engagement (%)")
    ax.set_ylim(0, 100); ax.grid(True, alpha=.25)
    fig.tight_layout(); fig.savefig(chart_path, dpi=140); plt.close(fig)

    # ---- document ----
    NAVY = colors.HexColor("#1E2761"); LIGHT = colors.HexColor("#F2F5FB")
    GREY = colors.HexColor("#5F6A85"); GOLD = colors.HexColor("#B07A1E")
    ss = getSampleStyleSheet()
    h1 = ParagraphStyle("h1", parent=ss["Title"], fontSize=20, textColor=NAVY,
                        alignment=0, spaceAfter=2)
    sub = ParagraphStyle("sub", parent=ss["Normal"], fontSize=10, textColor=GREY,
                         spaceAfter=10)
    h2 = ParagraphStyle("h2", parent=ss["Heading2"], fontSize=13, textColor=NAVY,
                        spaceBefore=12, spaceAfter=6)
    body = ParagraphStyle("body", parent=ss["Normal"], fontSize=10, leading=15,
                          spaceAfter=7)
    small = ParagraphStyle("small", parent=ss["Normal"], fontSize=8.5, leading=12,
                           textColor=GREY)

    doc = SimpleDocTemplate(pdf_path, pagesize=A4, leftMargin=20*mm, rightMargin=20*mm,
                            topMargin=18*mm, bottomMargin=18*mm)
    st = []
    st.append(Paragraph("Class Engagement Report", h1))
    st.append(Paragraph(
        f"{d.get('subject','Class')} &nbsp;|&nbsp; {d.get('class_name','')} "
        f"&nbsp;|&nbsp; {d.get('mode','recorded')} &nbsp;|&nbsp; {fmt(d['duration_s'])} "
        f"&nbsp;|&nbsp; generated {datetime.now().strftime('%d %B %Y, %H:%M')}", sub))

    rows = [["Average engagement", f"{d['average']:.0f}%"],
            ["Lowest point", f"{d['lowest_pct']:.0f}% at {fmt(d['lowest_time'])}"],
            ["Attention drops detected", str(len(d["events"]))],
            ["Measurement coverage", f"{d['coverage']:.0f}% of session "
                                     f"({d['gap_seconds']}s without a clear view)"]]
    if d.get("phone_share", 0) > 0:
        rows.append(["Phone use", f"visible for {d['phone_share']:.0f}% of session "
                                  f"(peak {d['phone_peak']} at once)"])
    tbl = Table(rows, colWidths=[58*mm, 97*mm])
    tbl.setStyle(TableStyle([
        ("BACKGROUND", (0,0), (-1,-1), LIGHT), ("TEXTCOLOR", (0,0), (0,-1), NAVY),
        ("FONTNAME", (0,0), (0,-1), "Helvetica-Bold"), ("FONTSIZE", (0,0), (-1,-1), 10),
        ("TOPPADDING", (0,0), (-1,-1), 7), ("BOTTOMPADDING", (0,0), (-1,-1), 7),
        ("LEFTPADDING", (0,0), (-1,-1), 10), ("GRID", (0,0), (-1,-1), .5, colors.white)]))
    st.append(tbl)

    st.append(Paragraph("Engagement across the session", h2))
    st.append(Image(chart_path, width=170*mm, height=66*mm))
    st.append(Paragraph(
        "The line traces class engagement over time; the dashed rule marks the session "
        "average. Shaded areas mark sustained periods below that average.", small))
    if d["coverage"] < 97:
        st.append(Spacer(1, 4))
        st.append(Paragraph(
            f"<b>Data quality note.</b> Students were clearly visible for "
            f"{d['coverage']:.0f}% of this session. The remaining {d['gap_seconds']} "
            f"seconds are excluded from the figures rather than counted as "
            f"disengagement, since an absence of measurement is not evidence of "
            f"inattention.", small))

    st.append(Paragraph("Moments that warrant attention", h2))
    if d["events"]:
        ev = [["Period", "Lowest engagement"]]
        for x in d["events"]:
            ev.append([f"{fmt(x['start'])} - {fmt(x['end'])}", f"{x['low']:.0f}%"])
        et = Table(ev, colWidths=[80*mm, 75*mm])
        et.setStyle(TableStyle([
            ("BACKGROUND", (0,0), (-1,0), NAVY), ("TEXTCOLOR", (0,0), (-1,0), colors.white),
            ("FONTNAME", (0,0), (-1,0), "Helvetica-Bold"),
            ("BACKGROUND", (0,1), (-1,-1), LIGHT), ("FONTSIZE", (0,0), (-1,-1), 9.5),
            ("TOPPADDING", (0,0), (-1,-1), 6), ("BOTTOMPADDING", (0,0), (-1,-1), 6),
            ("LEFTPADDING", (0,0), (-1,-1), 10), ("GRID", (0,0), (-1,-1), .5, colors.white)]))
        st.append(et)
    else:
        st.append(Paragraph("No sustained attention drops were detected.", body))

    z = d.get("zones") or []
    if z:
        st.append(Paragraph("Engagement by seating zone", h2))
        names = ["Front", "Middle", "Back"]
        zrows = [[""] + [f"Col {c+1}" for c in range(len(z[0]))]]
        for ri, row in enumerate(z):
            zrows.append([names[min(ri, 2)]] +
                         [("-" if v is None else f"{v:.0f}%") for v in row])
        zt = Table(zrows, colWidths=[30*mm] + [31*mm]*len(z[0]))
        zt.setStyle(TableStyle([
            ("BACKGROUND", (0,0), (-1,0), NAVY), ("TEXTCOLOR", (0,0), (-1,0), colors.white),
            ("BACKGROUND", (0,1), (0,-1), LIGHT), ("FONTSIZE", (0,0), (-1,-1), 9.5),
            ("ALIGN", (1,1), (-1,-1), "CENTER"),
            ("TOPPADDING", (0,0), (-1,-1), 6), ("BOTTOMPADDING", (0,0), (-1,-1), 6),
            ("GRID", (0,0), (-1,-1), .5, colors.HexColor("#DDDDDD"))]))
        st.append(zt)
        st.append(Paragraph("Rows are ordered from the teaching position toward the back "
                            "of the room.", small))

    st.append(Paragraph("Recommended actions", h2))
    st.append(Paragraph(
        "The following are derived from the patterns in this session's data. They are "
        "prompts for your own judgement rather than instructions: you know what was "
        "being taught at each moment, and that context takes precedence over the "
        "measurement.", body))
    for a in advice:
        st.append(Paragraph(f"<b>{a['title']}.</b> {a['body']}", body))

    st.append(Paragraph("How to use this report", h2))
    st.append(Paragraph(
        "This report is a private aid to your own reflection on the session. It is not an "
        "assessment of your teaching, and it is not a record of individual students: every "
        "figure describes the class as a whole, and the recording was deleted once "
        "analysis finished.", body))
    st.append(Paragraph(
        "The most useful way to read it is to work backwards from the flagged moments. "
        "Take each timestamp above and recall what was happening at that point - which "
        "concept you were introducing, which format you were using, how long you had been "
        "speaking without interruption. The system can identify <i>when</i> attention "
        "fell; only you can supply the <i>why</i>, and it is the combination that turns a "
        "measurement into a teaching decision.", body))
    st.append(Paragraph(
        "In the short term, the flagged periods indicate material worth revisiting next "
        "session. Over several sessions, comparing reports shows whether a change you made "
        "- different pacing, an added activity, an earlier break - had a measurable "
        "effect. And when preparing the material again, the periods of highest engagement "
        "are as informative as the lowest.", body))
    st.append(Paragraph(
        "A single report is one observation, not a conclusion. Engagement varies with time "
        "of day, point in term, and factors outside the classroom. Patterns persisting "
        "across sessions are meaningful; a single unusual session usually is not.", body))

    st.append(Paragraph("Interpreting the measurement", h2))
    st.append(Paragraph(
        "This system measures <b>behavioural</b> engagement: observable signals such as "
        "head orientation, posture, stillness and visible device use. It does not measure "
        "understanding or interest. A student may be attentive while still, or thinking "
        "deeply while looking away; apparent attention does not guarantee comprehension.", body))
    st.append(Paragraph(
        "Detection accuracy varies with conditions. Students furthest from the camera, "
        "partially occluded, or working head-down are measured less reliably, and fairness "
        "analysis has shown accuracy is not uniform across individuals. This is why results "
        "are reported only at class level and never used to characterise any individual "
        "student.", body))
    st.append(Spacer(1, 10))
    st.append(Paragraph(
        "Generated by LearnPulse. Source video is not retained; this report contains "
        "class-level aggregate results only.", small))

    doc.build(st)
    return pdf_path, safe


@app.get("/api/sessions/{sid}/report")
def download_report(sid: str, lp_token: str = Cookie(None)):
    u = require(lp_token)
    with db() as con:
        r = con.execute("SELECT subject,class_name FROM sessions WHERE id=? AND user_id=?",
                        (sid, u["id"])).fetchone()
    if not r:
        raise HTTPException(404, "Session not found")
    try:
        pdf, safe = build_pdf(sid, f"{r['subject']} {r['class_name']}")
    except HTTPException:
        raise
    except Exception as ex:
        import traceback; traceback.print_exc()
        raise HTTPException(500, f"Report generation failed: {ex}")
    return FileResponse(pdf, media_type="application/pdf",
                        filename=f"LearnPulse_Report_{safe}.pdf")


if __name__ == "__main__":
    print("\nLearnPulse ready at http://localhost:8000\n")
    uvicorn.run(app, host="127.0.0.1", port=8000)
